"""Document parsing: bytes + file type -> cleaned text + metadata.

Phase 2 only parses (no chunking/embeddings). Parser per format:
PDF→PyMuPDF, DOCX→python-docx, PPT(X)→python-pptx, CSV→pandas, HTML→BeautifulSoup,
TXT/MD/JSON→plain.
"""

from __future__ import annotations

import io
import json
from dataclasses import dataclass, field

SUPPORTED_TYPES = {"pdf", "docx", "txt", "md", "csv", "json", "ppt", "pptx", "html", "htm"}


@dataclass
class ParsedDocument:
    text: str
    metadata: dict = field(default_factory=dict)


class UnsupportedFileType(Exception):
    pass


def normalize_type(file_type: str) -> str:
    return file_type.lower().lstrip(".").strip()


def _parse_pdf(data: bytes) -> ParsedDocument:
    import fitz  # PyMuPDF

    with fitz.open(stream=data, filetype="pdf") as doc:
        text = "\n\n".join(page.get_text() for page in doc)
        return ParsedDocument(text=text, metadata={"page_count": doc.page_count})


def _parse_docx(data: bytes) -> ParsedDocument:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    return ParsedDocument(text=text, metadata={"paragraph_count": len(doc.paragraphs)})


def _parse_pptx(data: bytes) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return ParsedDocument(
        text="\n".join(parts), metadata={"slide_count": len(prs.slides._sldIdLst)}
    )


def _parse_csv(data: bytes) -> ParsedDocument:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(data))
    return ParsedDocument(
        text=df.to_csv(index=False),
        metadata={"row_count": int(df.shape[0]), "column_count": int(df.shape[1])},
    )


def _parse_json(data: bytes) -> ParsedDocument:
    parsed = json.loads(data.decode("utf-8", errors="replace"))
    return ParsedDocument(text=json.dumps(parsed, indent=2, ensure_ascii=False))


def _parse_html(data: bytes) -> ParsedDocument:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "lxml")
    return ParsedDocument(text=soup.get_text(separator="\n", strip=True))


def _parse_text(data: bytes) -> ParsedDocument:
    return ParsedDocument(text=data.decode("utf-8", errors="replace"))


def parse(data: bytes, file_type: str) -> ParsedDocument:
    ext = normalize_type(file_type)
    if ext not in SUPPORTED_TYPES:
        raise UnsupportedFileType(f"Unsupported file type: {file_type}")

    if ext == "pdf":
        result = _parse_pdf(data)
    elif ext == "docx":
        result = _parse_docx(data)
    elif ext in ("ppt", "pptx"):
        result = _parse_pptx(data)
    elif ext == "csv":
        result = _parse_csv(data)
    elif ext == "json":
        result = _parse_json(data)
    elif ext in ("html", "htm"):
        result = _parse_html(data)
    else:  # txt, md
        result = _parse_text(data)

    result.metadata["char_count"] = len(result.text)
    return result
